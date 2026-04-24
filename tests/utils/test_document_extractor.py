"""Tests for deeptutor.utils.document_extractor."""

from __future__ import annotations

import base64
import io

from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
import pytest

from deeptutor.utils.document_extractor import (
    CorruptDocumentError,
    DocumentTooLargeError,
    EmptyDocumentError,
    MAX_DOC_BYTES,
    MAX_EXTRACTED_CHARS_PER_DOC,
    UnsupportedDocumentError,
    extract_documents_from_records,
    extract_text_from_bytes,
    is_document_extension,
)


# ---------------------------------------------------------------------------
# Fixtures — generate office docs on the fly
# ---------------------------------------------------------------------------


def _make_docx(paragraphs: list[str]) -> bytes:
    doc = DocxDocument()
    for p in paragraphs:
        doc.add_paragraph(p)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx(sheets: dict[str, list[list[object]]]) -> bytes:
    wb = Workbook()
    default = wb.active
    first = True
    for name, rows in sheets.items():
        ws = default if first else wb.create_sheet()
        ws.title = name
        for row in rows:
            ws.append(row)
        first = False
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx(slides_text: list[list[str]]) -> bytes:
    prs = Presentation()
    for slide_texts in slides_text:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank-ish layout
        for i, text in enumerate(slide_texts):
            tb = slide.shapes.add_textbox(Inches(1), Inches(1 + i * 0.5), Inches(6), Inches(0.5))
            tb.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# is_document_extension
# ---------------------------------------------------------------------------


class TestIsDocumentExtension:
    def test_supported(self) -> None:
        assert is_document_extension("foo.pdf")
        assert is_document_extension("foo.DOCX")
        assert is_document_extension("report.xlsx")
        assert is_document_extension("deck.pptx")

    def test_unsupported(self) -> None:
        assert not is_document_extension("foo.png")
        assert not is_document_extension("foo.txt")
        assert not is_document_extension("foo")
        assert not is_document_extension("")


# ---------------------------------------------------------------------------
# extract_text_from_bytes — happy paths
# ---------------------------------------------------------------------------


class TestExtractDocx:
    def test_basic_paragraphs(self) -> None:
        data = _make_docx(["Hello world", "Second paragraph", ""])
        text = extract_text_from_bytes("doc.docx", data)
        assert "Hello world" in text
        assert "Second paragraph" in text


class TestExtractXlsx:
    def test_multiple_sheets(self) -> None:
        data = _make_xlsx(
            {
                "Alpha": [["a1", "b1"], ["a2", 42]],
                "Beta": [["x", "y"]],
            }
        )
        text = extract_text_from_bytes("book.xlsx", data)
        assert "--- Sheet: Alpha ---" in text
        assert "--- Sheet: Beta ---" in text
        assert "a1" in text and "42" in text
        assert "x" in text


class TestExtractPptx:
    def test_basic_slides(self) -> None:
        data = _make_pptx([["Slide 1 title", "Slide 1 body"], ["Slide 2 only text"]])
        text = extract_text_from_bytes("deck.pptx", data)
        assert "--- Slide 1 ---" in text
        assert "--- Slide 2 ---" in text
        assert "Slide 1 title" in text
        assert "Slide 2 only text" in text


class TestExtractPdf:
    def test_minimal_pdf(self) -> None:
        # Build a minimal valid PDF via pymupdf (dependency already in project).
        pytest.importorskip("fitz")
        import fitz  # noqa: WPS433

        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Hello PDF world")
        buf = io.BytesIO()
        doc.save(buf)
        doc.close()
        data = buf.getvalue()

        text = extract_text_from_bytes("sample.pdf", data)
        assert "Hello PDF world" in text
        assert "--- Page 1 ---" in text


# ---------------------------------------------------------------------------
# Failure modes
# ---------------------------------------------------------------------------


class TestFailureModes:
    def test_unsupported_extension(self) -> None:
        with pytest.raises(UnsupportedDocumentError):
            extract_text_from_bytes("foo.zip", b"\x00\x00")

    def test_empty_bytes(self) -> None:
        with pytest.raises(EmptyDocumentError):
            extract_text_from_bytes("foo.docx", b"")

    def test_too_large(self) -> None:
        fake = b"PK\x03\x04" + b"\x00" * (MAX_DOC_BYTES + 1)
        with pytest.raises(DocumentTooLargeError):
            extract_text_from_bytes("foo.docx", fake)

    def test_pdf_magic_mismatch(self) -> None:
        with pytest.raises(CorruptDocumentError):
            extract_text_from_bytes("foo.pdf", b"this is not a pdf")

    def test_ooxml_magic_mismatch(self) -> None:
        with pytest.raises(CorruptDocumentError):
            extract_text_from_bytes("foo.docx", b"not an office file")

    def test_corrupt_docx(self) -> None:
        # OOXML header but garbage body
        with pytest.raises(CorruptDocumentError):
            extract_text_from_bytes("foo.docx", b"PK\x03\x04" + b"\x00" * 512)

    def test_empty_docx_no_text(self) -> None:
        data = _make_docx([])  # no paragraphs
        with pytest.raises(EmptyDocumentError):
            extract_text_from_bytes("foo.docx", data)


# ---------------------------------------------------------------------------
# Truncation
# ---------------------------------------------------------------------------


class TestTruncation:
    def test_long_docx_is_truncated(self) -> None:
        # single paragraph of 250k chars → well over the 200k per-doc cap
        long_text = "a" * 250_000
        data = _make_docx([long_text])
        text = extract_text_from_bytes("big.docx", data)
        assert len(text) <= MAX_EXTRACTED_CHARS_PER_DOC + 200  # allow notice suffix
        assert "truncated" in text


# ---------------------------------------------------------------------------
# extract_documents_from_records
# ---------------------------------------------------------------------------


class TestExtractDocumentsFromRecords:
    def test_mixed_image_and_doc(self) -> None:
        docx_bytes = _make_docx(["hello there"])
        docx_b64 = base64.b64encode(docx_bytes).decode()
        image_b64 = base64.b64encode(b"\x89PNG\r\n\x1a\n").decode()

        records = [
            {"type": "image", "filename": "pic.png", "base64": image_b64, "mime_type": "image/png", "url": ""},
            {"type": "file", "filename": "note.docx", "base64": docx_b64, "mime_type": "", "url": ""},
        ]

        doc_texts, updated = extract_documents_from_records(records)

        assert len(doc_texts) == 1
        assert "[File: note.docx]" in doc_texts[0]
        assert "hello there" in doc_texts[0]

        # image record untouched
        assert updated[0]["base64"] == image_b64
        # doc record base64 cleared, extracted_chars set
        assert updated[1]["base64"] == ""
        assert updated[1]["extracted_chars"] > 0

    def test_unsupported_record_is_passthrough(self) -> None:
        records = [{"type": "file", "filename": "foo.zip", "base64": "AAAA", "mime_type": "", "url": ""}]
        doc_texts, updated = extract_documents_from_records(records)
        assert doc_texts == []
        assert updated[0]["base64"] == "AAAA"  # untouched — not a doc extension

    def test_failed_extraction_emits_error_marker(self) -> None:
        records = [{"type": "file", "filename": "bad.pdf", "base64": base64.b64encode(b"not a pdf").decode(), "mime_type": "", "url": ""}]
        doc_texts, updated = extract_documents_from_records(records)
        assert len(doc_texts) == 1
        assert "bad.pdf" in doc_texts[0]
        assert "could not be read" in doc_texts[0]
        assert updated[0]["base64"] == ""  # stripped even on failure

    def test_invalid_base64_emits_error_marker(self) -> None:
        records = [{"type": "file", "filename": "bad.docx", "base64": "!!!not base64!!!", "mime_type": "", "url": ""}]
        doc_texts, updated = extract_documents_from_records(records)
        # invalid base64 with validate=False may silently decode or emit error — both
        # paths end up as an error marker since resulting bytes won't pass magic check
        assert len(doc_texts) == 1
        assert "bad.docx" in doc_texts[0]
