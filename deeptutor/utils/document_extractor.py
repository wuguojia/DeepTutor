"""Document text extraction for chat attachments.

Bytes-in, text-out. Used by the chat turn runtime to inline the text of
user-dropped files into the ``effective_user_message`` sent to the LLM.

Two format families:
  * **Binary Office** (.pdf / .docx / .xlsx / .pptx) — parsed with pymupdf /
    python-docx / openpyxl / python-pptx.
  * **Text-like** (plain text, Markdown, source code, JSON, XML, CSV, …) —
    the extension set is imported from ``FileTypeRouter.TEXT_EXTENSIONS`` so
    the chat composer accepts every format the knowledge-base pipeline
    already ingests. Decoded with the same multi-encoding fallback chain.

Design mirrors ``nanobot/nanobot/utils/document.py`` but works on bytes
instead of file paths so the server never touches disk.
"""

from __future__ import annotations

import base64
from collections.abc import Iterable
import io
import logging
from pathlib import PurePosixPath

from deeptutor.services.rag.file_routing import FileTypeRouter

try:
    import fitz  # pymupdf
except ImportError:  # pragma: no cover
    fitz = None  # type: ignore[assignment]

try:
    from pypdf import PdfReader
    from pypdf.errors import FileNotDecryptedError as _PypdfNotDecryptedError
except ImportError:  # pragma: no cover
    PdfReader = None  # type: ignore[assignment]
    _PypdfNotDecryptedError = Exception  # type: ignore[assignment,misc]

try:
    from docx import Document as DocxDocument
except ImportError:  # pragma: no cover
    DocxDocument = None  # type: ignore[assignment]

try:
    from openpyxl import load_workbook
except ImportError:  # pragma: no cover
    load_workbook = None  # type: ignore[assignment]

try:
    from pptx import Presentation as PptxPresentation
except ImportError:  # pragma: no cover
    PptxPresentation = None  # type: ignore[assignment]

try:
    import ebooklib
    from ebooklib import epub
    import html as html_parser
except ImportError:  # pragma: no cover
    ebooklib = None  # type: ignore[assignment]
    epub = None  # type: ignore[assignment]
    html_parser = None  # type: ignore[assignment]


logger = logging.getLogger(__name__)


_OFFICE_EXTENSIONS: frozenset[str] = frozenset({".pdf", ".docx", ".xlsx", ".pptx", ".epub"})
# Text-like formats are sourced from the KB file router so chat and KB stay
# in sync. Adding a new code / config extension in one place propagates here.
TEXT_LIKE_EXTENSIONS: frozenset[str] = frozenset(FileTypeRouter.TEXT_EXTENSIONS)
SUPPORTED_DOC_EXTENSIONS: frozenset[str] = _OFFICE_EXTENSIONS | TEXT_LIKE_EXTENSIONS

MAX_DOC_BYTES = 10 * 1024 * 1024
MAX_TOTAL_DOC_BYTES = 25 * 1024 * 1024
MAX_EXTRACTED_CHARS_PER_DOC = 200_000
MAX_EXTRACTED_CHARS_TOTAL = 150_000

_PDF_MAGIC = b"%PDF-"
_OOXML_MAGIC = b"PK\x03\x04"
_EPUB_MAGIC = b"PK\x03\x04"  # EPUB is a ZIP file


class DocumentExtractionError(Exception):
    """Base class for extraction failures. ``str(exc)`` is user-friendly."""

    def __init__(self, message: str, filename: str = "") -> None:
        super().__init__(message)
        self.filename = filename


class UnsupportedDocumentError(DocumentExtractionError):
    pass


class CorruptDocumentError(DocumentExtractionError):
    pass


class EmptyDocumentError(DocumentExtractionError):
    pass


class DocumentTooLargeError(DocumentExtractionError):
    pass


def is_document_extension(filename: str) -> bool:
    return _ext(filename) in SUPPORTED_DOC_EXTENSIONS


def _ext(filename: str) -> str:
    return PurePosixPath(filename or "").suffix.lower()


def _truncate(text: str, max_length: int) -> str:
    if len(text) <= max_length:
        return text
    return text[:max_length] + f"... (truncated, {len(text)} chars total)"


def _check_magic(ext: str, data: bytes, filename: str) -> None:
    """Validate file header to catch extension spoofing.

    Only binary formats have well-known magic prefixes. Text-like extensions
    (code, markup, config, …) are decoded directly; a mislabeled binary blob
    either decodes as garbage or fails at decode time, which is fine.
    """
    if ext == ".pdf":
        if not data.startswith(_PDF_MAGIC):
            raise CorruptDocumentError(
                f"{filename} does not look like a PDF (bad header)", filename=filename
            )
    elif ext in {".docx", ".xlsx", ".pptx", ".epub"}:
        if not data.startswith(_OOXML_MAGIC):
            raise CorruptDocumentError(
                f"{filename} does not look like a valid Office/EPUB file (bad header)",
                filename=filename,
            )


def extract_text_from_bytes(filename: str, data: bytes) -> str:
    """Extract text from a single document's raw bytes.

    Raises a ``DocumentExtractionError`` subclass on failure. Successful
    output is truncated to ``MAX_EXTRACTED_CHARS_PER_DOC`` with a notice.
    """
    if not data:
        raise EmptyDocumentError(f"{filename} is empty", filename=filename)
    if len(data) > MAX_DOC_BYTES:
        raise DocumentTooLargeError(
            f"{filename} exceeds the {MAX_DOC_BYTES // (1024 * 1024)} MB per-file limit",
            filename=filename,
        )

    ext = _ext(filename)
    if ext not in SUPPORTED_DOC_EXTENSIONS:
        raise UnsupportedDocumentError(
            f"{filename} has unsupported extension '{ext}'", filename=filename
        )

    _check_magic(ext, data, filename)

    if ext == ".pdf":
        text = _extract_pdf(data, filename)
    elif ext == ".docx":
        text = _extract_docx(data, filename)
    elif ext == ".xlsx":
        text = _extract_xlsx(data, filename)
    elif ext == ".pptx":
        text = _extract_pptx(data, filename)
    elif ext == ".epub":
        text = _extract_epub(data, filename)
    elif ext in TEXT_LIKE_EXTENSIONS:
        text = _extract_text_like(data, filename)
    else:  # pragma: no cover - guarded above
        raise UnsupportedDocumentError(f"{filename}: unreachable", filename=filename)

    if not text.strip():
        raise EmptyDocumentError(f"{filename}: no extractable text", filename=filename)
    return _truncate(text, MAX_EXTRACTED_CHARS_PER_DOC)


def _extract_pdf(data: bytes, filename: str) -> str:
    if fitz is not None:
        try:
            with fitz.open(stream=data, filetype="pdf") as doc:
                if doc.is_encrypted and not doc.authenticate(""):
                    raise CorruptDocumentError(
                        f"{filename} is encrypted and cannot be read", filename=filename
                    )
                pages = [f"--- Page {i} ---\n{page.get_text() or ''}" for i, page in enumerate(doc, 1)]
            return "\n\n".join(pages)
        except CorruptDocumentError:
            raise
        except Exception as exc:
            logger.warning("pymupdf failed on %s: %s — falling back to pypdf", filename, exc)

    if PdfReader is None:
        raise CorruptDocumentError(
            f"{filename}: no PDF reader available (install pymupdf or pypdf)",
            filename=filename,
        )
    try:
        reader = PdfReader(io.BytesIO(data))
        if getattr(reader, "is_encrypted", False):
            raise CorruptDocumentError(
                f"{filename} is encrypted and cannot be read", filename=filename
            )
        pages = [f"--- Page {i} ---\n{page.extract_text() or ''}" for i, page in enumerate(reader.pages, 1)]
        return "\n\n".join(pages)
    except CorruptDocumentError:
        raise
    except _PypdfNotDecryptedError as exc:
        raise CorruptDocumentError(
            f"{filename} is encrypted and cannot be read", filename=filename
        ) from exc
    except Exception as exc:
        raise CorruptDocumentError(f"{filename}: failed to read PDF ({exc})", filename=filename) from exc


def _extract_docx(data: bytes, filename: str) -> str:
    if DocxDocument is None:
        raise CorruptDocumentError(
            f"{filename}: python-docx not installed", filename=filename
        )
    try:
        doc = DocxDocument(io.BytesIO(data))
    except Exception as exc:
        raise CorruptDocumentError(
            f"{filename}: failed to open DOCX ({exc})", filename=filename
        ) from exc
    paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n\n".join(paragraphs)


def _extract_xlsx(data: bytes, filename: str) -> str:
    if load_workbook is None:
        raise CorruptDocumentError(
            f"{filename}: openpyxl not installed", filename=filename
        )
    try:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as exc:
        raise CorruptDocumentError(
            f"{filename}: failed to open XLSX ({exc})", filename=filename
        ) from exc
    try:
        sheets: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                sheets.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
        return "\n\n".join(sheets)
    finally:
        wb.close()


def _extract_pptx(data: bytes, filename: str) -> str:
    if PptxPresentation is None:
        raise CorruptDocumentError(
            f"{filename}: python-pptx not installed", filename=filename
        )
    try:
        prs = PptxPresentation(io.BytesIO(data))
    except Exception as exc:
        raise CorruptDocumentError(
            f"{filename}: failed to open PPTX ({exc})", filename=filename
        ) from exc
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            _collect_pptx_shape_text(shape, slide_text)
        if slide_text:
            slides.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))
    return "\n\n".join(slides)


def _extract_epub(data: bytes, filename: str) -> str:
    """Extract text from EPUB (electronic book) files.

    EPUB is a ZIP-based format containing HTML/XHTML content.
    This function extracts text from all document items in reading order,
    including metadata, chapter structure, and embedded images descriptions.
    """
    if epub is None:
        raise CorruptDocumentError(
            f"{filename}: ebooklib not installed", filename=filename
        )
    try:
        book = epub.read_epub(io.BytesIO(data))
    except Exception as exc:
        raise CorruptDocumentError(
            f"{filename}: failed to open EPUB ({exc})", filename=filename
        ) from exc

    # Extract metadata
    metadata_parts: list[str] = []
    title = book.get_metadata('DC', 'title')
    if title:
        metadata_parts.append(f"Title: {title[0][0]}")

    author = book.get_metadata('DC', 'creator')
    if author:
        metadata_parts.append(f"Author: {author[0][0]}")

    language = book.get_metadata('DC', 'language')
    if language:
        metadata_parts.append(f"Language: {language[0][0]}")

    publisher = book.get_metadata('DC', 'publisher')
    if publisher:
        metadata_parts.append(f"Publisher: {publisher[0][0]}")

    # Add metadata header if available
    result_parts: list[str] = []
    if metadata_parts:
        result_parts.append("=== EPUB Metadata ===\n" + "\n".join(metadata_parts))

    # Extract chapters with improved structure
    chapters: list[str] = []
    image_count = 0

    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            try:
                content = item.get_content().decode('utf-8', errors='ignore')
                # Extract chapter title if available
                chapter_title = _extract_chapter_title(content)
                # Strip HTML tags to get plain text with enhanced processing
                text = _strip_html_tags_enhanced(content)
                if text.strip():
                    if chapter_title:
                        chapters.append(f"--- {chapter_title} ---\n{text}")
                    else:
                        chapters.append(text)
            except Exception as exc:
                logger.warning(f"Failed to extract chapter from {filename}: {exc}")
                continue
        elif item.get_type() == ebooklib.ITEM_IMAGE:
            # Track images for metadata
            image_count += 1

    if chapters:
        result_parts.append("\n\n".join(chapters))

    # Add image summary if images exist
    if image_count > 0:
        result_parts.append(f"\n[This EPUB contains {image_count} embedded images]")

    return "\n\n".join(result_parts)


def _strip_html_tags(html_content: str) -> str:
    """Strip HTML tags from content and return plain text.

    Uses html.parser to unescape entities and simple regex to remove tags.
    """
    import re

    # Remove HTML tags
    text = re.sub(r'<script[^>]*>.*?</script>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<[^>]+>', ' ', text)

    # Unescape HTML entities
    if html_parser:
        text = html_parser.unescape(text)

    # Clean up whitespace
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'\n\s*\n', '\n\n', text)

    return text.strip()


def _extract_chapter_title(html_content: str) -> str:
    """Extract chapter title from HTML content.

    Looks for h1, h2, or title tags at the beginning of the content.
    """
    import re

    # Try h1 first
    h1_match = re.search(r'<h1[^>]*>(.*?)</h1>', html_content, re.IGNORECASE | re.DOTALL)
    if h1_match:
        title = re.sub(r'<[^>]+>', '', h1_match.group(1))
        if html_parser:
            title = html_parser.unescape(title)
        return title.strip()

    # Try h2
    h2_match = re.search(r'<h2[^>]*>(.*?)</h2>', html_content, re.IGNORECASE | re.DOTALL)
    if h2_match:
        title = re.sub(r'<[^>]+>', '', h2_match.group(1))
        if html_parser:
            title = html_parser.unescape(title)
        return title.strip()

    # Try title tag
    title_match = re.search(r'<title[^>]*>(.*?)</title>', html_content, re.IGNORECASE | re.DOTALL)
    if title_match:
        title = re.sub(r'<[^>]+>', '', title_match.group(1))
        if html_parser:
            title = html_parser.unescape(title)
        return title.strip()

    return ""


def _strip_html_tags_enhanced(html_content: str) -> str:
    """Enhanced HTML tag stripping with better structure preservation.

    Preserves code blocks, lists, tables, and other structured content.
    """
    import re

    # Remove script and style tags
    text = re.sub(r'<script[^>]*>.*?</script>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)

    # Preserve code blocks with markers
    text = re.sub(r'<pre[^>]*>(.*?)</pre>', r'\n[CODE]\n\1\n[/CODE]\n', text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<code[^>]*>(.*?)</code>', r'`\1`', text, flags=re.DOTALL | re.IGNORECASE)

    # Convert lists to readable format
    text = re.sub(r'<li[^>]*>', '\n• ', text, flags=re.IGNORECASE)
    text = re.sub(r'</li>', '', text, flags=re.IGNORECASE)
    text = re.sub(r'</?[ou]l[^>]*>', '\n', text, flags=re.IGNORECASE)

    # Convert tables to tab-separated format
    text = re.sub(r'<tr[^>]*>', '\n', text, flags=re.IGNORECASE)
    text = re.sub(r'</tr>', '', text, flags=re.IGNORECASE)
    text = re.sub(r'<t[hd][^>]*>', '', text, flags=re.IGNORECASE)
    text = re.sub(r'</t[hd]>', '\t', text, flags=re.IGNORECASE)
    text = re.sub(r'</?table[^>]*>', '\n[TABLE]\n', text, flags=re.IGNORECASE)

    # Add line breaks for block elements
    text = re.sub(r'<br\s*/?>', '\n', text, flags=re.IGNORECASE)
    text = re.sub(r'</?p[^>]*>', '\n', text, flags=re.IGNORECASE)
    text = re.sub(r'</?div[^>]*>', '\n', text, flags=re.IGNORECASE)
    text = re.sub(r'</h[1-6]>', '\n', text, flags=re.IGNORECASE)

    # Remove remaining HTML tags
    text = re.sub(r'<[^>]+>', ' ', text)

    # Unescape HTML entities
    if html_parser:
        text = html_parser.unescape(text)

    # Clean up whitespace but preserve structure
    text = re.sub(r' +', ' ', text)
    text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)

    return text.strip()


def _extract_text_like(data: bytes, filename: str) -> str:
    """Decode a plain-text / code / config / markup file.

    Uses the same encoding fallback chain as the KB pipeline
    (``FileTypeRouter.decode_bytes``) so a GBK-encoded Python file or a
    UTF-8-BOM Markdown works the same way in both places.
    """
    try:
        return FileTypeRouter.decode_bytes(data)
    except Exception as exc:  # pragma: no cover - decode_bytes never raises
        raise CorruptDocumentError(
            f"{filename}: failed to decode text ({exc})", filename=filename
        ) from exc


def _collect_pptx_shape_text(shape, out: list[str]) -> None:
    """Recurse into groups + tables, same semantics as nanobot's version."""
    sub_shapes = getattr(shape, "shapes", None)
    if sub_shapes is not None:
        for sub in sub_shapes:
            _collect_pptx_shape_text(sub, out)
        return

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            line = "\t".join(cell for cell in cells if cell)
            if line:
                out.append(line)
        return

    text = getattr(shape, "text", "")
    if text:
        out.append(text)


def extract_documents_from_records(
    records: Iterable[dict],
) -> tuple[list[str], list[dict]]:
    """Process a list of attachment records from the WS payload.

    Parameters
    ----------
    records:
        Raw attachment records as parsed by the turn runtime
        (``{"type", "url", "base64", "filename", "mime_type"}``).

    Returns
    -------
    (doc_texts, updated_records)
        ``doc_texts`` is a list of strings formatted as
        ``"[File: <name>]\\n<text>"`` (one per processed or skipped doc).
        ``updated_records`` is the input list with the ``base64`` field
        cleared on successfully-extracted docs (to save DB space), an
        ``extracted_chars`` field added, and the extracted plain text
        stored under ``extracted_text`` so the chat UI can preview office
        documents without re-running the parser. Image / non-document
        records are returned unchanged.
    """
    doc_texts: list[str] = []
    updated: list[dict] = []
    total_bytes = 0
    total_chars = 0
    over_quota = False

    for raw in records:
        record = dict(raw)
        filename = str(record.get("filename") or "")
        if not is_document_extension(filename):
            updated.append(record)
            continue

        b64 = record.get("base64") or ""
        if not b64:
            updated.append(record)
            continue

        if over_quota:
            doc_texts.append(
                f"[File: {filename} — skipped: total attachment quota exceeded]"
            )
            record["base64"] = ""
            record["extracted_chars"] = 0
            updated.append(record)
            continue

        try:
            data = base64.b64decode(b64, validate=False)
        except Exception as exc:
            doc_texts.append(f"[File: {filename} — could not be read: invalid base64 ({exc})]")
            record["base64"] = ""
            record["extracted_chars"] = 0
            updated.append(record)
            continue

        if total_bytes + len(data) > MAX_TOTAL_DOC_BYTES:
            over_quota = True
            doc_texts.append(
                f"[File: {filename} — skipped: total attachment quota exceeded]"
            )
            record["base64"] = ""
            record["extracted_chars"] = 0
            updated.append(record)
            continue

        total_bytes += len(data)

        try:
            text = extract_text_from_bytes(filename, data)
        except DocumentExtractionError as exc:
            logger.info("Document extraction failed for %s: %s", filename, exc)
            doc_texts.append(f"[File: {filename} — could not be read: {exc}]")
            record["base64"] = ""
            record["extracted_chars"] = 0
            updated.append(record)
            continue

        remaining_budget = MAX_EXTRACTED_CHARS_TOTAL - total_chars
        if remaining_budget <= 0:
            doc_texts.append(
                f"[File: {filename} — skipped: total extracted-text quota exceeded]"
            )
            record["base64"] = ""
            record["extracted_chars"] = 0
            updated.append(record)
            continue

        if len(text) > remaining_budget:
            text = (
                text[:remaining_budget]
                + f"... (truncated, {len(text)} chars total; turn quota hit)"
            )

        total_chars += len(text)
        doc_texts.append(f"[File: {filename}]\n{text}")
        record["base64"] = ""
        record["extracted_chars"] = len(text)
        record["extracted_text"] = text
        updated.append(record)

    return doc_texts, updated
